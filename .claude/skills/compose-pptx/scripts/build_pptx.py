"""Build a .pptx from slide-deck YAML by filling duplicated template samples.

Stage 3 of the narrative -> YAML -> pptx pipeline (the compose-pptx skill).

Inputs:
  --yaml      slide-outline.yaml (schema: narrative-to-slide-outline/references/
              slide_yaml_schema.md). Its embedded data: blocks are authoritative.
  --template  a "sample slide deck" .pptx: each slide is a styled example layout.
  --mapping   mapping.json: [{"yaml_index": int, "template_index": int,
              "notes": str}], produced by the matching step (Step 3).
  --out       output .pptx (default: <yaml dir>/deck.pptx).

What it does, per YAML slide:
  1. Duplicate the mapped template sample slide (deep-copy its shape tree into a
     fresh slide created from the same layout, re-relating images/etc.).
  2. Fill it: title text, body (light Markdown), tables (cell rewrite, rows
     grown/shrunk), pictures (replaced from image path), charts, speaker notes.
  3. After all slides are built, delete the original template sample slides.

CHART DESIGN NOTE (deliberate): a chart is NOT reused from the sample via
chart.replace_data(). Duplicating a slide makes the copy share the *same* chart
part as the original; replace_data() then mutates a shared part and corrupts
sibling slides that map to the same sample. So instead we drop the copied chart
graphic-frame and add a fresh native chart (add_chart) at the same position.
Each output slide thus owns an independent chart that still inherits the deck's
theme colors/fonts. (Cost: a sample chart's bespoke styling is not preserved.)

Run with an interpreter where python-pptx imports (here: /opt/uv-venv/bin/python).
"""

import argparse
import copy
import json
import re
import sys
from pathlib import Path

import yaml
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx.oxml.ns import qn
from pptx.util import Inches

PLACEHOLDER_RE = re.compile(r"\{\{[^}]+\}\}")
BOLD_RE = re.compile(r"\*\*(.+?)\*\*")

CHART_TYPE_MAP = {
    "bar_chart": XL_CHART_TYPE.COLUMN_CLUSTERED,  # vertical columns = the usual slide "bar"
    "line_chart": XL_CHART_TYPE.LINE,
    "histogram": XL_CHART_TYPE.COLUMN_CLUSTERED,
}

# Default chart geometry when the matched sample has no chart frame to inherit.
DEFAULT_CHART_GEOM = (Inches(1.0), Inches(1.8), Inches(8.0), Inches(4.8))
# Default table geometry when the matched sample has no table to reuse.
DEFAULT_TABLE_GEOM = (Inches(0.8), Inches(2.2), Inches(8.4), Inches(4.2))


# --------------------------------------------------------------------------- #
# Slide duplication
# --------------------------------------------------------------------------- #
def duplicate_slide(prs, src_slide):
    """Append a deep copy of src_slide and return the new slide.

    Copies the whole shape tree, then re-relates non-layout relationships
    (images, media, hyperlinks). Chart parts are intentionally NOT re-related;
    callers replace charts with fresh ones (see module docstring). References in
    the copied XML are remapped to whatever new rIds the destination assigns.
    """
    dest = prs.slides.add_slide(src_slide.slide_layout)

    # Drop the placeholders the layout seeded into the new slide.
    for shp in list(dest.shapes):
        shp._element.getparent().remove(shp._element)

    # Copy every shape element from the source.
    for shp in src_slide.shapes:
        dest.shapes._spTree.append(copy.deepcopy(shp._element))

    # Re-relate parts and remap rId references in the copied XML.
    rId_map = {}
    for rId, rel in src_slide.part.rels.items():
        if rel.reltype == RT.SLIDE_LAYOUT:
            continue  # dest already relates to its own layout
        if rel.reltype == RT.CHART:
            continue  # charts are rebuilt fresh, not re-related
        if rel.is_external:
            new_rId = dest.part.rels.get_or_add_ext_rel(rel.reltype, rel.target_ref)
        else:
            new_rId = dest.part.relate_to(rel.target_part, rel.reltype)
        rId_map[rId] = new_rId

    _remap_rels(dest.shapes._spTree, rId_map)
    return dest


def _remap_rels(spTree, rId_map):
    r_ns = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
    prefix = "{" + r_ns + "}"
    for el in spTree.iter():
        for name, val in list(el.attrib.items()):
            if name.startswith(prefix) and val in rId_map:
                el.set(name, rId_map[val])


# --------------------------------------------------------------------------- #
# Shape lookup helpers
# --------------------------------------------------------------------------- #
def _is_title(shape):
    try:
        return shape.is_placeholder and shape.placeholder_format.type in (
            PP_PLACEHOLDER.TITLE,
            PP_PLACEHOLDER.CENTER_TITLE,
        )
    except Exception:
        return False


def find_title(slide):
    try:
        if slide.shapes.title is not None:
            return slide.shapes.title
    except Exception:
        pass
    for shape in slide.shapes:
        if _is_title(shape):
            return shape
    return None


def find_body(slide, title_shape):
    """First non-title text placeholder; else first non-title text box.

    Excludes the title by underlying element, not Python identity: python-pptx
    builds a fresh proxy object on every shape access, so `shape is title_shape`
    would never match the same element. _is_title is also checked as a guard so
    body text is never written into a title placeholder.
    """
    title_el = title_shape._element if title_shape is not None else None
    candidates = []
    for shape in slide.shapes:
        if title_el is not None and shape._element is title_el:
            continue
        if _is_title(shape) or not getattr(shape, "has_text_frame", False):
            continue
        if getattr(shape, "has_chart", False) or getattr(shape, "has_table", False):
            continue
        is_ph = bool(getattr(shape, "is_placeholder", False))
        candidates.append((0 if is_ph else 1, shape))
    candidates.sort(key=lambda t: t[0])
    return candidates[0][1] if candidates else None


def _geom(shape):
    return (shape.left, shape.top, shape.width, shape.height)


# --------------------------------------------------------------------------- #
# Content fillers
# --------------------------------------------------------------------------- #
def _fmt(v):
    if isinstance(v, bool):
        return str(v)
    if isinstance(v, int):
        return str(v)
    if isinstance(v, float):
        return ("%g" % v)
    return str(v)


def _add_runs(paragraph, text):
    """Add runs to a paragraph, honoring **bold** spans."""
    pos = 0
    for m in BOLD_RE.finditer(text):
        if m.start() > pos:
            run = paragraph.add_run()
            run.text = text[pos:m.start()]
        run = paragraph.add_run()
        run.text = m.group(1)
        run.font.bold = True
        pos = m.end()
    if pos < len(text):
        run = paragraph.add_run()
        run.text = text[pos:]


def render_body(text_frame, body):
    """Render a light-Markdown body into a text frame.

    Supports: blank-line paragraph breaks, '-'/'*' bullets (indent = nesting
    level), numbered lines (kept verbatim), and **bold** inline. {{token}}
    placeholders are stripped. Preserves the placeholder's own bullet styling by
    only setting paragraph text/level.
    """
    text_frame.clear()
    body = PLACEHOLDER_RE.sub("", body or "")
    lines = [ln for ln in body.split("\n")]

    # Trim leading/trailing blank lines.
    while lines and not lines[0].strip():
        lines.pop(0)
    while lines and not lines[-1].strip():
        lines.pop()

    if not lines:
        return

    first = True
    for raw in lines:
        if not raw.strip():
            continue  # collapse blank lines; paragraphs are per content line
        indent = len(raw) - len(raw.lstrip(" "))
        content = raw.strip()
        level = 0
        m = re.match(r"^[-*]\s+(.*)$", content)
        if m:
            content = m.group(1)
            level = min(indent // 2, 4)
        if first:
            p = text_frame.paragraphs[0]
            first = False
        else:
            p = text_frame.add_paragraph()
        if level:
            p.level = level
        _add_runs(p, content)


def _set_table_dims(table, nrows, ncols):
    """Reshape a table to nrows x ncols by cloning/removing <a:gridCol>/<a:tr>/<a:tc>.

    Columns are adjusted first (on existing rows + the grid), then rows are
    grown by cloning the last <a:tr>, which by then already has the right
    column count. Cloned cells keep the sample's styling; fill_table overwrites
    their text.
    """
    tbl = table._tbl
    grid = tbl.find(qn("a:tblGrid"))
    cols = grid.findall(qn("a:gridCol"))
    cur_c = len(cols)
    # Preserve the sample table's total width so a grown table still fits.
    orig_total_w = sum(int(c.get("w")) for c in cols if c.get("w"))
    if ncols > cur_c:
        for _ in range(ncols - cur_c):
            grid.append(copy.deepcopy(cols[-1]))
        for tr in tbl.findall(qn("a:tr")):
            tcs = tr.findall(qn("a:tc"))
            for _ in range(ncols - cur_c):
                tr.append(copy.deepcopy(tcs[-1]))
    elif ncols < cur_c:
        for gc in cols[ncols:]:
            grid.remove(gc)
        for tr in tbl.findall(qn("a:tr")):
            for tc in tr.findall(qn("a:tc"))[ncols:]:
                tr.remove(tc)
    if ncols != cur_c and orig_total_w:
        # Redistribute the original total width equally across the new columns.
        each = orig_total_w // ncols
        for gc in grid.findall(qn("a:gridCol")):
            gc.set("w", str(each))

    trs = tbl.findall(qn("a:tr"))
    cur_r = len(trs)
    if nrows > cur_r:
        for _ in range(nrows - cur_r):
            tbl.append(copy.deepcopy(trs[-1]))
    elif nrows < cur_r:
        for tr in trs[nrows:]:
            tbl.remove(tr)


def fill_table(graphic_frame, entry):
    table = graphic_frame.table
    headers = entry.get("headers")
    rows = entry.get("rows", []) or []
    grid = ([headers] if headers else []) + rows
    if not grid:
        return
    ncols = max(len(r) for r in grid)
    _set_table_dims(table, len(grid), ncols)
    for r, row_vals in enumerate(grid):
        for c in range(ncols):
            val = row_vals[c] if c < len(row_vals) else ""
            table.cell(r, c).text = _fmt(val)


def replace_picture(slide, pic_shape, image_path, caption):
    left, top, width, height = _geom(pic_shape)
    pic_shape._element.getparent().remove(pic_shape._element)
    new_pic = slide.shapes.add_picture(str(image_path), left, top, width, height)
    if caption:
        new_pic._element.nvPicPr.cNvPr.set("descr", str(caption))
    return new_pic


def add_native_chart(slide, geom, entry):
    ctype = entry["type"]
    chart_data = CategoryChartData()
    if ctype == "histogram":
        chart_data.categories = entry.get("bins", [])
        chart_data.add_series(
            entry.get("y_axis", "Frequency"),
            tuple(entry.get("frequencies", [])),
        )
    else:
        chart_data.categories = entry.get("categories", [])
        for series in entry.get("series", []):
            chart_data.add_series(series["name"], tuple(series["values"]))
    left, top, width, height = geom
    gframe = slide.shapes.add_chart(
        CHART_TYPE_MAP[ctype], left, top, width, height, chart_data
    )
    chart = gframe.chart
    n_series = 1 if ctype == "histogram" else len(entry.get("series", []))
    chart.has_legend = n_series > 1
    if chart.has_legend:
        chart.legend.include_in_layout = False
    return chart


# --------------------------------------------------------------------------- #
# Per-slide assembly
# --------------------------------------------------------------------------- #
def build_slide(prs, src_slide, spec, warnings):
    dest = duplicate_slide(prs, src_slide)

    # Classify the copied shapes (snapshot before we mutate the tree).
    chart_frames = [s for s in dest.shapes if getattr(s, "has_chart", False)]
    table_frames = [s for s in dest.shapes if getattr(s, "has_table", False)]
    pictures = list(s for s in dest.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE)

    # Charts: record geometry, then remove the copied (shared-part) frames.
    chart_geoms = [_geom(f) for f in chart_frames]
    for f in chart_frames:
        f._element.getparent().remove(f._element)

    # Title: prefer a proper title placeholder; fall back to the largest text box
    # (by area). This handles BLANK-layout templates where all text lives in plain
    # text boxes with no placeholder type.
    title_shape = find_title(dest)
    if title_shape is None:
        text_boxes = [
            s for s in dest.shapes
            if getattr(s, "has_text_frame", False)
            and not getattr(s, "has_chart", False)
            and not getattr(s, "has_table", False)
        ]
        if text_boxes:
            title_shape = max(text_boxes, key=lambda s: s.width * s.height)
    if title_shape is not None:
        title_shape.text_frame.clear()
        _add_runs(title_shape.text_frame.paragraphs[0], spec.get("title", ""))

    # Body text.
    body_shape = find_body(dest, title_shape)
    if body_shape is not None:
        render_body(body_shape.text_frame, spec.get("body", ""))

    # Clear any remaining text frames so leftover template sample text does not
    # bleed through into the output slide.
    title_el = title_shape._element if title_shape is not None else None
    body_el = body_shape._element if body_shape is not None else None
    for shape in dest.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue
        if title_el is not None and shape._element is title_el:
            continue
        if body_el is not None and shape._element is body_el:
            continue
        shape.text_frame.clear()

    # Data entries, matched by type to the sample's objects (in order).
    data = spec.get("data") or {}
    chart_i = table_i = pic_i = 0
    for name, entry in data.items():
        etype = entry.get("type")
        if etype in CHART_TYPE_MAP:
            if chart_i < len(chart_geoms):
                geom = chart_geoms[chart_i]
            elif pictures:
                # No native chart frame — steal the first picture's position.
                geom = _geom(pictures[0])
                pictures[0]._element.getparent().remove(pictures[0]._element)
                pictures.pop(0)
                warnings.append(
                    f"slide '{spec.get('title', '')}': no chart frame in template "
                    f"sample; replaced picture with chart '{name}'"
                )
            else:
                geom = DEFAULT_CHART_GEOM
                warnings.append(
                    f"slide '{spec.get('title', '')}': no chart frame in template "
                    f"sample; placed chart '{name}' at default position"
                )
            add_native_chart(dest, geom, entry)
            chart_i += 1
        elif etype == "table":
            if table_i < len(table_frames):
                fill_table(table_frames[table_i], entry)
                table_i += 1
            else:
                # No table in sample — add a fresh one at the default position.
                headers = entry.get("headers")
                rows = entry.get("rows", []) or []
                grid = ([headers] if headers else []) + rows
                nrows = len(grid) if grid else 1
                ncols = max(len(r) for r in grid) if grid else 1
                left, top, width, height = DEFAULT_TABLE_GEOM
                gframe = dest.shapes.add_table(nrows, ncols, left, top, width, height)
                fill_table(gframe, entry)
                warnings.append(
                    f"slide '{spec.get('title', '')}': no table in template sample "
                    f"for '{name}'; added at default position"
                )
                table_i += 1
        elif etype == "image":
            path = entry.get("path")
            if path and not Path(path).exists():
                warnings.append(f"image not found: {path}")
            if pic_i < len(pictures) and path and Path(path).exists():
                replace_picture(dest, pictures[pic_i], path, entry.get("caption"))
                pic_i += 1
            elif path and Path(path).exists():
                left, top, width, height = DEFAULT_CHART_GEOM
                pic = dest.shapes.add_picture(str(path), left, top, width, height)
                if entry.get("caption"):
                    pic._element.nvPicPr.cNvPr.set("descr", str(entry["caption"]))

    # Speaker notes.
    notes = spec.get("speaker_notes")
    if notes:
        dest.notes_slide.notes_text_frame.text = notes

    return dest


def remove_original_slides(prs, original_count):
    """Un-list the first `original_count` slides (the template samples).

    Only the <p:sldId> entry is removed from the slide-id list, which is what
    determines display order. The slide parts themselves become orphaned (no
    sldId references them) and are simply ignored by PowerPoint -- the same
    state the app itself leaves behind when you delete a slide.
    """
    sldIdLst = prs.slides._sldIdLst
    for sldId in list(sldIdLst)[:original_count]:
        sldIdLst.remove(sldId)


# --------------------------------------------------------------------------- #
# Main
# --------------------------------------------------------------------------- #
def main():
    parser = argparse.ArgumentParser(description="Build .pptx from slide YAML + template.")
    parser.add_argument("--yaml", required=True)
    parser.add_argument("--template", required=True)
    parser.add_argument("--mapping", required=True)
    parser.add_argument("--out", default=None)
    args = parser.parse_args()

    yaml_path = Path(args.yaml)
    deck = yaml.safe_load(yaml_path.read_text(encoding="utf-8"))
    mapping = json.loads(Path(args.mapping).read_text(encoding="utf-8"))
    out_path = Path(args.out) if args.out else yaml_path.parent / "deck.pptx"

    slides_spec = deck.get("slides", [])
    map_by_yaml = {m["yaml_index"]: m for m in mapping}

    prs = Presentation(args.template)
    originals = list(prs.slides)
    original_count = len(originals)
    if original_count == 0:
        print("Error: template has no sample slides", file=sys.stderr)
        sys.exit(1)

    # Core properties from YAML metadata.
    if deck.get("title"):
        prs.core_properties.title = deck["title"]
    if deck.get("author"):
        prs.core_properties.author = deck["author"]

    warnings = []
    for i, spec in enumerate(slides_spec):
        m = map_by_yaml.get(i)
        t_idx = m["template_index"] if m else 0
        if t_idx < 0 or t_idx >= original_count:
            warnings.append(f"slide {i}: template_index {t_idx} out of range; used 0")
            t_idx = 0
        build_slide(prs, originals[t_idx], spec, warnings)

    remove_original_slides(prs, original_count)

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))

    print(f"Wrote {out_path} ({len(slides_spec)} slides)")
    if warnings:
        print(f"\n{len(warnings)} warning(s):")
        for w in warnings:
            print(f"  - {w}")


if __name__ == "__main__":
    main()
