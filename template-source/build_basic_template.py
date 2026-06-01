"""Generate a basic, neutral/minimal sample-slide TEMPLATE deck for compose-pptx-2.

The point of this template is NOT the dummy content — it is to ship one *styled
sample* of every layout compose-pptx-2 knows how to reuse, with **native**
PowerPoint objects (real charts, a real table), so the skill can duplicate a
sample and refill it:

  0  Title            10  Closing
  1  Bullets
  2  Section divider  (empty body -> divider sample)
  3  Two-column
  4  Bar chart        (native clustered column)
  5  Line chart       (native line w/ markers)
  6  Pie chart        (native pie)
  7  Table            (native a:tbl)
  8  Card grid        (3 styled tiles -> clone_grid points at one)
  9  Image            (a real picture, so set_image can swap it)

Charts are left on the deck's THEME accents (python-pptx writes schemeClr
references, not hard-coded colours), so a later theme swap recolours them.

Run:  /opt/uv-venv/bin/python build_basic_template.py [out.pptx]
"""

import sys
from pathlib import Path

from PIL import Image, ImageDraw
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.oxml.ns import qn
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_LINE_DASH_STYLE, MSO_THEME_COLOR
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION

# ---- Palette: Apple-style design tokens ------------------------------------
# Single near-black ink + a SINGLE interactive accent (Action Blue). White /
# parchment canvases punctuated by near-black "tiles" whose colour change is the
# section divider. No second accent, no decorative colour.
TEXT     = RGBColor(0x1D, 0x1D, 0x1F)   # ink — every headline & body line
MUTED    = RGBColor(0x6E, 0x6E, 0x73)   # secondary text (Apple gray)
ACCENT   = RGBColor(0x00, 0x66, 0xCC)   # Action Blue — the ONE accent
CARD     = RGBColor(0xF5, 0xF5, 0xF7)   # parchment tile fill
HAIR     = RGBColor(0xE0, 0xE0, 0xE0)   # hairline borders / rules
GRID     = RGBColor(0xD9, 0xD9, 0xD9)   # faint grey value-axis gridlines
DIVBG    = RGBColor(0x1D, 0x1D, 0x1F)   # near-black divider/closing tile
ONDARK   = RGBColor(0xFF, 0xFF, 0xFF)   # text on dark tiles
DARKMUTED = RGBColor(0xCC, 0xCC, 0xCC)  # secondary text on dark tiles

# Chart accent palette written into the theme's accent1..6. python-pptx charts
# carry NO explicit series colour, so the renderer auto-assigns accent1→series 1,
# accent2→series 2, … — and compose-pptx's add_chart/set_chart_data inherit the
# same theme accents. So this ordered list is the chart colour candidate pool;
# accent1 is series-1 / first pie slice. Apple-system hues, calmed for slides.
CHART_ACCENTS = ["0066CC",   # accent1 — Action Blue   (series 1 / first slice)
                 "E8893B",   # accent2 — orange        (max contrast w/ blue)
                 "34A853",   # accent3 — green
                 "C44D8A",   # accent4 — pink
                 "5AC8DC",   # accent5 — teal
                 "5E5CE6"]   # accent6 — indigo

# Soften the saturated accents slightly: render every coloured chart mark (bars,
# lines, pie slices) and the two accent rules (title underline + takeaway bar) at
# 87% opacity so the blues/oranges read a touch lighter against white. OOXML alpha
# is in thousandths of a percent, so 87% == "87000".
#
# NOTE: an <a:alpha> placed in a THEME accent *definition* is NOT honoured when a
# chart series resolves it through a <a:schemeClr> reference (verified: LibreOffice
# renders such bars at full saturation). So opacity must live as a transform on the
# reference itself — i.e. on the per-series/per-point schemeClr we write below — not
# on the theme colour. Using schemeClr (not srgbClr) keeps charts theme-swappable.
ACCENT_ALPHA = "87000"
ACCENT_THEME = [MSO_THEME_COLOR.ACCENT_1, MSO_THEME_COLOR.ACCENT_2,
                MSO_THEME_COLOR.ACCENT_3, MSO_THEME_COLOR.ACCENT_4,
                MSO_THEME_COLOR.ACCENT_5, MSO_THEME_COLOR.ACCENT_6]

# Cross-platform sans: Arial ships on Windows+Mac and substitutes to a sans
# (Arimo) on ChromeOS/web, so no font embedding is needed and the file stays tiny.
FONT = "Arial"

EMU_IN = 914400
SW, SH = 13.333, 7.5                  # 16:9 wide, inches
MARGIN = 0.9


def _box(slide, x, y, w, h):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    return tb, tf


def _kill_shadow(shape):
    """python-pptx's `shape.shadow.inherit = False` does NOT emit an explicit
    empty <a:effectLst/>, so PowerPoint and LibreOffice still paint the preset
    drop shadow. Inject an empty effectLst to truly disable it. (effectLst sits
    after <a:ln> in CT_ShapeProperties; these shapes carry no scene3d/extLst, so
    appending keeps the schema order valid.)"""
    spPr = shape._element.spPr
    if spPr.find(qn("a:effectLst")) is None:
        spPr.append(spPr.makeelement(qn("a:effectLst"), {}))
    # Autoshapes also carry a <p:style> whose <a:effectRef idx="N"> points at a
    # THEME effect style that includes an outer shadow; LibreOffice honours that
    # ref even when spPr has an empty effectLst. Force the ref to idx 0 (none).
    style = shape._element.find(qn("p:style"))
    if style is not None:
        eref = style.find(qn("a:effectRef"))
        if eref is not None:
            eref.set("idx", "0")


def _fill_alpha(shape, alpha=ACCENT_ALPHA):
    """Render a solid-filled shape at partial opacity. python-pptx has no opacity
    API, so inject an <a:alpha> (thousandths of a percent) into the fill colour's
    <a:srgbClr>. Call AFTER the solid fore_color has been set."""
    srgbClr = shape._element.spPr.find(qn("a:solidFill")).find(qn("a:srgbClr"))
    srgbClr.append(srgbClr.makeelement(qn("a:alpha"), {"val": alpha}))


def _style_gridlines(chart):
    """Replace a cartesian chart's default solid-black value-axis gridlines (the
    horizontal scale rules) with a faint light-grey DOTTED rule. Pie charts have
    no value axis, so this is only called for the column/line charts."""
    va = chart.value_axis
    va.has_major_gridlines = True
    ln = va.major_gridlines.format.line
    ln.color.rgb = GRID
    ln.width = Pt(0.75)
    ln.dash_style = MSO_LINE_DASH_STYLE.ROUND_DOT


def _alpha_on(scheme_clr):
    """Tag an <a:schemeClr> element with ACCENT_ALPHA opacity."""
    scheme_clr.append(scheme_clr.makeelement(qn("a:alpha"), {"val": ACCENT_ALPHA}))


def _series_fill_alpha(chart):
    """Column/bar/histogram charts: give each series an explicit theme-accent SOLID
    FILL at ACCENT_ALPHA opacity. Written as a schemeClr reference (+ <a:alpha>), so
    it reproduces the renderer's auto accent1..N order yet stays theme-swappable."""
    for i, ser in enumerate(chart.series):
        ser.format.fill.solid()
        ser.format.fill.fore_color.theme_color = ACCENT_THEME[i % 6]
        spPr = ser.format._element.get_or_add_spPr()
        _alpha_on(spPr.find(qn("a:solidFill")).find(qn("a:schemeClr")))


def _line_alpha(chart):
    """Line charts: the visible mark is the series LINE, so colour the line with a
    theme accent at ACCENT_ALPHA opacity (schemeClr + <a:alpha>)."""
    for i, ser in enumerate(chart.series):
        ser.format.line.color.theme_color = ACCENT_THEME[i % 6]
        spPr = ser.format._element.get_or_add_spPr()
        ln = spPr.find(qn("a:ln"))
        _alpha_on(ln.find(qn("a:solidFill")).find(qn("a:schemeClr")))


def _pie_alpha(chart):
    """Pie charts: each slice is a separate data point, so fill every point with its
    auto accent (accent1.. for slice 1..) at ACCENT_ALPHA opacity."""
    for i, pt in enumerate(chart.plots[0].series[0].points):
        pt.format.fill.solid()
        pt.format.fill.fore_color.theme_color = ACCENT_THEME[i % 6]
        spPr = pt.format._element.get_or_add_spPr()
        _alpha_on(spPr.find(qn("a:solidFill")).find(qn("a:schemeClr")))


def _style_run(run, size, color=TEXT, bold=False, italic=False):
    f = run.font
    f.size = Pt(size)
    f.bold = bold
    f.italic = italic
    f.color.rgb = color
    f.name = FONT


def add_title(slide, text, sub=None):
    _, tf = _box(slide, MARGIN, 0.55, SW - 2 * MARGIN, 1.0)
    p = tf.paragraphs[0]
    _style_run(p.add_run(), 30, TEXT, bold=True)
    p.runs[0].text = text
    # thin accent rule beneath the title (kept close under the title text)
    rule = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(MARGIN), Inches(1.13), Inches(5.2), Pt(3)
    )
    rule.fill.solid(); rule.fill.fore_color.rgb = ACCENT; _fill_alpha(rule)
    rule.line.fill.background()
    rule.shadow.inherit = False; _kill_shadow(rule)
    if sub:
        _, sf = _box(slide, MARGIN, 1.55, SW - 2 * MARGIN, 0.6)
        sp = sf.paragraphs[0]
        _style_run(sp.add_run(), 17, MUTED)
        sp.runs[0].text = sub


def number_paragraph(p):
    """Turn a paragraph into an auto-numbered list item (1. 2. 3. ...)."""
    pPr = p._p.get_or_add_pPr()
    # marL = hanging width = the number->text gap. The "N." glyph is wider than a
    # bullet, so it needs a larger marL than bullets to leave the same visible gap.
    pPr.set("marL", "247650")
    pPr.set("indent", "-247650")
    for tag in ("a:buNone", "a:buAutoNum", "a:buChar"):
        ex = pPr.find(qn(tag))
        if ex is not None:
            pPr.remove(ex)
    bu = pPr.makeelement(qn("a:buAutoNum"), {"type": "arabicPeriod"})
    pPr.append(bu)


def bullet_paragraph(p):
    """Turn a paragraph into a plain bulleted list item (• ...)."""
    pPr = p._p.get_or_add_pPr()
    # marL = hanging width = the bullet->text gap. A touch of breathing room.
    pPr.set("marL", "203200")
    pPr.set("indent", "-203200")
    for tag in ("a:buNone", "a:buAutoNum", "a:buChar"):
        ex = pPr.find(qn(tag))
        if ex is not None:
            pPr.remove(ex)
    pPr.append(pPr.makeelement(qn("a:buFont"), {"typeface": FONT}))
    pPr.append(pPr.makeelement(qn("a:buChar"), {"char": "•"}))


def set_legend(chart, pos):
    """Set a chart legend with a guaranteed `val` on <c:legendPos> (python-pptx
    1.0.2 writes an empty <c:legendPos/> for some positions, which makes
    PowerPoint refuse to render the chart). pos in {'r','l','t','b'}."""
    chart.has_legend = True
    leg = chart.legend._element
    lp = leg.find(qn("c:legendPos"))
    if lp is None:
        lp = leg.makeelement(qn("c:legendPos"), {})
        leg.insert(0, lp)
    lp.set("val", pos)
    chart.legend.include_in_layout = False


def add_insight(slide, text):
    """A one-line takeaway placed above the content (between the title rule and
    the object). compose-pptx-2 refills this as the slide's body text."""
    _, tf = _box(slide, MARGIN, 1.6, SW - 2 * MARGIN, 0.6)
    p = tf.paragraphs[0]
    _style_run(p.add_run(), 17, MUTED); p.runs[0].text = text


def add_caption(slide, text, y):
    """A normal-text line below the content — same format as the intro/insight
    line (17pt, muted, non-italic), so it reads as ordinary prose, not a caption."""
    _, tf = _box(slide, MARGIN, y, SW - 2 * MARGIN, 0.6)
    p = tf.paragraphs[0]
    _style_run(p.add_run(), 17, MUTED); p.runs[0].text = text


def blank(prs):
    # slide_layouts[6] == "Blank" in the default python-pptx template
    return prs.slides.add_slide(prs.slide_layouts[6])


def slide_title(prs):
    s = blank(prs)
    # Near-black title tile (mirrors the divider/closing surface).
    s.background.fill.solid(); s.background.fill.fore_color.rgb = DIVBG
    _, tf = _box(s, MARGIN, 2.7, SW - 2 * MARGIN, 1.4)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    _style_run(p.add_run(), 44, ONDARK, bold=True); p.runs[0].text = "Presentation Title"
    _, sf = _box(s, MARGIN, 4.0, SW - 2 * MARGIN, 0.9)
    sp = sf.paragraphs[0]
    _style_run(sp.add_run(), 20, DARKMUTED); sp.runs[0].text = "Subtitle / author / date"
    return s


def slide_bullets(prs):
    s = blank(prs); add_title(s, "Numbered List")
    _, tf = _box(s, MARGIN, 1.9, SW - 2 * MARGIN, 4.8)
    samples = ["First key point goes here", "Second supporting point",
               "Third point with a bit more detail", "Fourth point"]
    for i, t in enumerate(samples):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(10); p.line_spacing = 1.15
        _style_run(p.add_run(), 18, TEXT); p.runs[0].text = t
        number_paragraph(p)
    return s


def slide_bullets_plain(prs):
    s = blank(prs); add_title(s, "Bulleted List")
    _, tf = _box(s, MARGIN, 1.9, SW - 2 * MARGIN, 4.8)
    samples = ["First key point goes here", "Second supporting point",
               "Third point with a bit more detail", "Fourth point"]
    for i, t in enumerate(samples):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(10); p.line_spacing = 1.15
        _style_run(p.add_run(), 18, TEXT); p.runs[0].text = t
        bullet_paragraph(p)
    return s


def slide_divider(prs):
    s = blank(prs)
    # Near-black full-bleed "tile": the surface colour change IS the divider.
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = DIVBG
    _, tf = _box(s, MARGIN, 3.0, SW - 2 * MARGIN, 1.5)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    _style_run(p.add_run(), 34, ONDARK, bold=True); p.runs[0].text = "Section Divider"
    return s


def slide_two_col(prs):
    s = blank(prs); add_title(s, "Two Columns")
    colw = (SW - 2 * MARGIN - 0.6) / 2
    for i, (head, items) in enumerate([
        ("Left heading", ["Left point one", "Left point two", "Left point three"]),
        ("Right heading", ["Right point one", "Right point two", "Right point three"]),
    ]):
        x = MARGIN + i * (colw + 0.6)
        _, tf = _box(s, x, 1.9, colw, 4.6)
        hp = tf.paragraphs[0]; hp.space_after = Pt(18)
        _style_run(hp.add_run(), 20, ACCENT, bold=True); hp.runs[0].text = head
        for t in items:
            p = tf.add_paragraph(); p.space_after = Pt(8)
            _style_run(p.add_run(), 17, TEXT); p.runs[0].text = t
    return s


def _columns(s, n, fontsz):
    gap = 0.5
    colw = (SW - 2 * MARGIN - (n - 1) * gap) / n
    for i in range(n):
        x = MARGIN + i * (colw + gap)
        _, tf = _box(s, x, 1.9, colw, 4.6)
        hp = tf.paragraphs[0]; hp.space_after = Pt(16)
        _style_run(hp.add_run(), 18, ACCENT, bold=True); hp.runs[0].text = f"Heading {i+1}"
        for j in range(3):
            p = tf.add_paragraph(); p.space_after = Pt(8)
            _style_run(p.add_run(), fontsz, TEXT); p.runs[0].text = f"Point {j+1}"


def slide_three_col(prs):
    s = blank(prs); add_title(s, "Three Columns")
    _columns(s, 3, 16)
    return s


def slide_four_col(prs):
    s = blank(prs); add_title(s, "Four Columns")
    _columns(s, 4, 14)
    return s


def _chart_frame(slide):
    # Lowered top leaves room for the insight line above the chart.
    top = 2.35
    return (Inches(MARGIN), Inches(top),
            Inches(SW - 2 * MARGIN), Inches(SH - top - 0.6))


def slide_bar(prs):
    s = blank(prs); add_title(s, "Bar Chart")
    add_insight(s, "One-line takeaway about what this chart shows.")
    cd = CategoryChartData()
    cd.categories = ["Q1", "Q2", "Q3", "Q4"]
    cd.add_series("Plan", (18, 22, 25, 30))
    cd.add_series("Actual", (16, 24, 23, 33))
    gf = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, *_chart_frame(s), cd)
    plot = gf.chart.plots[0]
    plot.overlap = -50      # negative -> a clear gap between the two bars in each group
    plot.gap_width = 160    # wider category gap -> slightly narrower bars
    _style_gridlines(gf.chart)
    _series_fill_alpha(gf.chart)
    set_legend(gf.chart, "b")
    return s


def slide_line(prs):
    s = blank(prs); add_title(s, "Line Chart")
    add_insight(s, "One-line takeaway about the trend over time.")
    cd = CategoryChartData()
    cd.categories = ["Jan", "Feb", "Mar", "Apr", "May", "Jun"]
    cd.add_series("Revenue", (12, 15, 14, 19, 22, 27))
    gf = s.shapes.add_chart(XL_CHART_TYPE.LINE_MARKERS, *_chart_frame(s), cd)
    _style_gridlines(gf.chart)
    _line_alpha(gf.chart)
    set_legend(gf.chart, "b")
    return s


def slide_pie(prs):
    s = blank(prs); add_title(s, "Pie Chart")
    add_insight(s, "One-line takeaway about the composition.")
    cd = CategoryChartData()
    cd.categories = ["Enterprise", "SMB", "Consumer", "Other"]
    cd.add_series("Share", (42, 28, 22, 8))
    # A pie reads better narrow: shrink the frame width to 60% of the full content
    # width and centre it horizontally (the cartesian charts keep full width).
    fx, fy, fw, fh = _chart_frame(s)
    narrow = Emu(int(fw * 0.6))
    cx = Emu(int((prs.slide_width - narrow) / 2))
    gf = s.shapes.add_chart(XL_CHART_TYPE.PIE, cx, fy, narrow, fh, cd)
    _pie_alpha(gf.chart)
    set_legend(gf.chart, "r")
    gf.chart.plots[0].has_data_labels = True
    return s


def slide_histogram(prs):
    """Distribution of a single variable over pre-binned intervals. Built as a
    single-series clustered column with near-zero gap so bars sit touching — the
    textbook histogram look. compose-pptx maps `type: histogram` here."""
    s = blank(prs); add_title(s, "Histogram")
    add_insight(s, "One-line takeaway about the distribution.")
    cd = CategoryChartData()
    cd.categories = ["0-20", "20-40", "40-60", "60-80", "80-100"]
    cd.add_series("Count", (3, 12, 45, 30, 10))
    gf = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, *_chart_frame(s), cd)
    plot = gf.chart.plots[0]
    plot.gap_width = 10     # near-touching bars -> histogram look
    _style_gridlines(gf.chart)
    _series_fill_alpha(gf.chart)
    gf.chart.has_legend = False
    return s


def slide_split_chart(prs):
    """Split layout: commentary text on the left, native chart on the right."""
    s = blank(prs); add_title(s, "Text + Chart")
    # Left: heading + bullets
    _, tf = _box(s, MARGIN, 2.0, 4.7, 4.6)
    hp = tf.paragraphs[0]; hp.space_after = Pt(28)
    _style_run(hp.add_run(), 20, ACCENT, bold=True); hp.runs[0].text = "Key takeaway heading"
    for t in ["Supporting point one", "Supporting point two",
              "Supporting point three"]:
        p = tf.add_paragraph(); p.space_after = Pt(10)
        _style_run(p.add_run(), 17, TEXT); p.runs[0].text = t
    # Right: chart
    rx = 6.1
    cd = CategoryChartData()
    cd.categories = ["Q1", "Q2", "Q3", "Q4"]
    cd.add_series("Actual", (16, 24, 23, 33))
    gf = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED,
                            Inches(rx), Inches(2.0),
                            Inches(SW - MARGIN - rx), Inches(4.4), cd)
    _style_gridlines(gf.chart)
    _series_fill_alpha(gf.chart)
    gf.chart.has_legend = False
    return s


# "No Style, No Grid" built-in table style GUID — strips python-pptx's default
# coloured/banded style so we can paint a clean hairline table ourselves.
_TBL_NO_STYLE = "{2D5ABB26-0587-4C30-8999-92F81FD0307C}"


def _table_plain(tbl):
    """Drop the default banded/filled table style and turn off the auto first-row
    and banding flags, leaving an unstyled grid to paint by hand."""
    tbl.first_row = False; tbl.horz_banding = False
    tblPr = tbl._tbl.find(qn("a:tblPr"))
    if tblPr is not None:
        sid = tblPr.find(qn("a:tableStyleId"))
        if sid is None:
            sid = tblPr.makeelement(qn("a:tableStyleId"), {})
            tblPr.append(sid)
        sid.text = _TBL_NO_STYLE


def _cell_bottom_border(cell, color, emu):
    """Add only a bottom hairline to a cell (lines precede <a:fill> in
    CT_TableCellProperties, so insert at the front of tcPr)."""
    tcPr = cell._tc.get_or_add_tcPr()
    ex = tcPr.find(qn("a:lnB"))
    if ex is not None:
        tcPr.remove(ex)
    ln = tcPr.makeelement(qn("a:lnB"),
                          {"w": str(emu), "cap": "flat", "cmpd": "sng", "algn": "ctr"})
    sf = ln.makeelement(qn("a:solidFill"), {})
    sf.append(ln.makeelement(qn("a:srgbClr"), {"val": str(color)}))
    ln.append(sf)
    tcPr.insert(0, ln)


def _populate_table(tbl, data, cols, rows, fontsz=14):
    """Minimal Apple-style table: white cells, ink text, left-aligned labels and
    right-aligned numerics, a strong ink rule under the header, and thin hairline
    row separators. No heavy header fill, no banding, no shadow."""
    _table_plain(tbl)
    for r in range(rows):
        header = (r == 0)
        for c in range(cols):
            cell = tbl.cell(r, c)
            cell.text = str(data[r][c])
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.fill.solid(); cell.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            cell.margin_left = Inches(0.12); cell.margin_right = Inches(0.12)
            para = cell.text_frame.paragraphs[0]
            _style_run(para.runs[0], fontsz, TEXT, bold=header)
            para.alignment = PP_ALIGN.LEFT if c == 0 else PP_ALIGN.RIGHT
            if header:
                _cell_bottom_border(cell, TEXT, 19050)   # ~1.5pt ink under header
            elif r < rows - 1:
                _cell_bottom_border(cell, HAIR, 9525)     # ~1px hairline separator


def slide_two_tables(prs):
    """Two tables side by side (fill_table --table-index 0 / 1 fills both)."""
    s = blank(prs); add_title(s, "Two Tables")
    add_insight(s, "Intro line describing the two tables.")
    gap = 0.5
    tw = (SW - 2 * MARGIN - gap) / 2
    data = [["Segment", "Customers", "Revenue"], ["Enterprise", "45", "112"],
            ["SMB", "620", "84"], ["Consumer", "5,300", "39"]]
    for i in range(2):
        x = MARGIN + i * (tw + gap)
        gf = s.shapes.add_table(4, 3, Inches(x), Inches(2.3), Inches(tw), Inches(2.9))
        _populate_table(gf.table, data, 3, 4)
    add_caption(s, "Source / footnote line below the tables.", 5.5)
    return s


def slide_table_chart(prs):
    """Split: a table on the left, its chart on the right."""
    s = blank(prs); add_title(s, "Table + Chart")
    add_insight(s, "Intro line describing the table and its chart.")
    tw = 5.4
    data = [["Quarter", "Plan", "Actual"], ["Q1", "18", "16"], ["Q2", "22", "24"],
            ["Q3", "25", "23"], ["Q4", "30", "33"]]
    gf = s.shapes.add_table(5, 3, Inches(MARGIN), Inches(2.3), Inches(tw), Inches(3.2))
    _populate_table(gf.table, data, 3, 5)
    rx = MARGIN + tw + 0.6
    cd = CategoryChartData()
    cd.categories = ["Q1", "Q2", "Q3", "Q4"]
    cd.add_series("Plan", (18, 22, 25, 30))
    cd.add_series("Actual", (16, 24, 23, 33))
    g2 = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(rx), Inches(2.3),
                            Inches(SW - MARGIN - rx), Inches(3.2), cd)
    plot = g2.chart.plots[0]; plot.overlap = -50; plot.gap_width = 160
    _style_gridlines(g2.chart)
    _series_fill_alpha(g2.chart)
    set_legend(g2.chart, "b")
    add_caption(s, "Source / footnote line below the table and chart.", 5.7)
    return s


def slide_table(prs):
    s = blank(prs); add_title(s, "Table")
    add_insight(s, "Intro line describing what the table shows.")
    rows, cols = 4, 3
    gf = s.shapes.add_table(rows, cols, Inches(MARGIN), Inches(2.3),
                            Inches(SW - 2 * MARGIN), Inches(2.9))
    tbl = gf.table
    data = [["Segment", "Customers", "Revenue"],
            ["Enterprise", "45", "112"],
            ["SMB", "620", "84"],
            ["Consumer", "5,300", "39"]]
    _populate_table(tbl, data, cols, rows, fontsz=15)
    add_caption(s, "Source / footnote line below the table.", 5.5)
    return s


def slide_cards(prs):
    s = blank(prs); add_title(s, "Card Grid")
    add_insight(s, "Intro line above the cards.")
    n = 3
    gap = 0.4
    cw = (SW - 2 * MARGIN - (n - 1) * gap) / n
    cy, ch = 2.35, 2.7
    for i in range(n):
        x = MARGIN + i * (cw + gap)
        card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  Inches(x), Inches(cy), Inches(cw), Inches(ch))
        card.adjustments[0] = 0.07   # ~18px (rounded.lg) corner, not pill-round
        card.fill.solid(); card.fill.fore_color.rgb = CARD
        card.line.color.rgb = HAIR; card.line.width = Pt(0.75)
        card.shadow.inherit = False; _kill_shadow(card)
        tf = card.text_frame; tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        hp = tf.paragraphs[0]; hp.alignment = PP_ALIGN.CENTER
        _style_run(hp.add_run(), 30, ACCENT, bold=True); hp.runs[0].text = f"{(i+1)*12}%"
        bp = tf.add_paragraph(); bp.alignment = PP_ALIGN.CENTER
        _style_run(bp.add_run(), 15, MUTED); bp.runs[0].text = f"Metric {i+1} label"
    add_caption(s, "Note / source line below the cards.", 5.35)
    return s


def slide_image(prs, img_path):
    s = blank(prs); add_title(s, "Image")
    s.shapes.add_picture(str(img_path), Inches(MARGIN), Inches(2.0),
                         Inches(SW - 2 * MARGIN), Inches(4.3))
    return s


def slide_image_split(prs, img_path):
    """Split layout: commentary text on the left, image on the right."""
    s = blank(prs); add_title(s, "Text + Image")
    _, tf = _box(s, MARGIN, 2.0, 4.7, 4.6)
    hp = tf.paragraphs[0]; hp.space_after = Pt(18)
    _style_run(hp.add_run(), 20, ACCENT, bold=True); hp.runs[0].text = "Heading"
    for t in ["Supporting point one", "Supporting point two",
              "Supporting point three"]:
        p = tf.add_paragraph(); p.space_after = Pt(10)
        _style_run(p.add_run(), 17, TEXT); p.runs[0].text = t
    rx = 6.1
    s.shapes.add_picture(str(img_path), Inches(rx), Inches(2.0),
                         Inches(SW - MARGIN - rx), Inches(4.2))
    return s


def slide_conclusion(prs):
    """Summary slide before the closing: a few takeaways + a highlighted box."""
    s = blank(prs); add_title(s, "Conclusion")
    _, tf = _box(s, MARGIN, 1.9, SW - 2 * MARGIN, 3.0)
    for i, t in enumerate(["Summary point one", "Summary point two",
                           "Summary point three"]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(12)
        _style_run(p.add_run(), 18, TEXT); p.runs[0].text = t
    # highlighted key-takeaway box with an accent left bar
    bw, bh, by = SW - 2 * MARGIN, 1.1, 5.4
    box = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(MARGIN), Inches(by),
                             Inches(bw), Inches(bh))
    box.fill.solid(); box.fill.fore_color.rgb = CARD
    box.line.fill.background(); box.shadow.inherit = False; _kill_shadow(box)
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(MARGIN), Inches(by),
                             Inches(0.09), Inches(bh))
    bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT; _fill_alpha(bar)
    bar.line.fill.background(); bar.shadow.inherit = False; _kill_shadow(bar)
    tb = box.text_frame
    tb.word_wrap = True
    tb.vertical_anchor = MSO_ANCHOR.MIDDLE
    tb.margin_left = Inches(0.3)
    p = tb.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    _style_run(p.add_run(), 18, TEXT, bold=True)
    p.runs[0].text = "Key message goes here."
    return s


def slide_closing(prs):
    s = blank(prs)
    # Near-black tile to bookend the deck (mirrors the divider surface).
    s.background.fill.solid(); s.background.fill.fore_color.rgb = DIVBG
    _, tf = _box(s, MARGIN, 3.0, SW - 2 * MARGIN, 1.4)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    _style_run(p.add_run(), 40, ONDARK, bold=True); p.runs[0].text = "Thank You"
    _, sf = _box(s, MARGIN, 4.2, SW - 2 * MARGIN, 0.7)
    sp = sf.paragraphs[0]; sp.alignment = PP_ALIGN.CENTER
    _style_run(sp.add_run(), 18, DARKMUTED); sp.runs[0].text = "name@example.com"
    return s


def make_placeholder_png(path):
    w, h = 1280, 720
    img = Image.new("RGB", (w, h), (0xEC, 0xEE, 0xF1))
    d = ImageDraw.Draw(img)
    d.rectangle([2, 2, w - 3, h - 3], outline=(0xCF, 0xD4, 0xDA), width=3)
    d.line([0, 0, w, h], fill=(0xD7, 0xDB, 0xE0), width=2)
    d.line([0, h, w, 0], fill=(0xD7, 0xDB, 0xE0), width=2)
    img.save(path)


def patch_theme_accents(pptx_path, accents):
    """Rewrite accent1..6 in every theme part of the saved .pptx (python-pptx has
    no theme-colour API, so we post-process the package ZIP). Native charts carry
    no explicit colour, so they auto-colour from these accents — and so will any
    chart compose-pptx later builds from this template. `accents` is six 'RRGGBB'
    strings, accent1 first."""
    import zipfile, re, os
    src = zipfile.ZipFile(pptx_path)
    infos = src.infolist()
    blobs = {i.filename: src.read(i.filename) for i in infos}
    src.close()
    theme_re = re.compile(r"ppt/theme/theme\d+\.xml$")
    for name in blobs:
        if not theme_re.match(name):
            continue
        xml = blobs[name].decode("utf-8")
        for idx, hexv in enumerate(accents, start=1):
            # Replace the single colour child (srgbClr or sysClr) inside <a:accentN>.
            xml = re.sub(
                r'(<a:accent%d>)\s*<a:(?:srgbClr|sysClr)\b[^>]*/>\s*(</a:accent%d>)' % (idx, idx),
                r'\g<1><a:srgbClr val="%s"/>\g<2>' % hexv, xml, count=1)
        blobs[name] = xml.encode("utf-8")
    tmp = str(pptx_path) + ".tmp"
    with zipfile.ZipFile(tmp, "w", zipfile.ZIP_DEFLATED) as out:
        for i in infos:                       # preserve original entry order
            out.writestr(i, blobs[i.filename])
    os.replace(tmp, str(pptx_path))


def main():
    out = Path(sys.argv[1]) if len(sys.argv) > 1 else Path(__file__).with_name("basic_template.pptx")
    img_path = out.with_name("_placeholder.png")
    make_placeholder_png(img_path)

    prs = Presentation()
    prs.slide_width = Emu(int(SW * EMU_IN))
    prs.slide_height = Emu(int(SH * EMU_IN))
    # python-pptx leaves the inherited type="screen4x3" on <p:sldSz>; the real
    # dimensions are 16:9, so correct the label to match.
    prs._element.find(qn("p:sldSz")).set("type", "screen16x9")

    slide_title(prs)
    slide_bullets(prs)
    slide_bullets_plain(prs)
    slide_divider(prs)
    slide_two_col(prs)
    slide_three_col(prs)
    slide_four_col(prs)
    slide_bar(prs)
    slide_line(prs)
    slide_pie(prs)
    slide_histogram(prs)
    slide_split_chart(prs)
    slide_table(prs)
    slide_two_tables(prs)
    slide_table_chart(prs)
    slide_cards(prs)
    slide_image(prs, img_path)
    slide_image_split(prs, img_path)
    slide_conclusion(prs)
    slide_closing(prs)

    # Seed every slide with a notes slide so set_notes.py (which clones an
    # existing notesSlide) has a master+sample to work from; without this,
    # speaker_notes from the YAML are silently dropped.
    for slide in prs.slides:
        slide.notes_slide.notes_text_frame.text = "Speaker notes here."

    # python-pptx adds the notesMaster PART + relationship but does NOT register
    # it in presentation.xml via <p:notesMasterIdLst>. PowerPoint tolerates the
    # orphan relationship, but stricter consumers (macOS Quick Look / Keynote)
    # reject the whole deck ("no displayable content"). Wire it up explicitly.
    R_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
    pres = prs._element
    if pres.find(qn("p:notesMasterIdLst")) is None:
        notes_rid = next((rel.rId for rel in prs.part.rels.values()
                          if rel.reltype.endswith("notesMaster")), None)
        if notes_rid:
            nmidlst = pres.makeelement(qn("p:notesMasterIdLst"), {})
            nmid = pres.makeelement(qn("p:notesMasterId"), {qn("r:id"): notes_rid})
            nmidlst.append(nmid)
            pres.find(qn("p:sldMasterIdLst")).addnext(nmidlst)

    prs.save(str(out))
    # Recolour the theme accents AFTER saving (post-process the package ZIP).
    patch_theme_accents(out, CHART_ACCENTS)
    print(f"Wrote {out} with {len(prs.slides._sldIdLst)} slides.")


if __name__ == "__main__":
    main()
